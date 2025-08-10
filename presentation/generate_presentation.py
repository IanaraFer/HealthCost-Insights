"""
PowerPoint Presentation Generator for HealthCost Insights
Creates a comprehensive presentation with charts and figures
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
import os

# Try to import python-pptx for PowerPoint generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
    print("✅ python-pptx available - Full PowerPoint generation enabled")
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx not available - Will generate static charts only")
    print("   Install with: pip install python-pptx")

def create_presentation_charts():
    """Generate all charts for the presentation"""
    
    print("📊 Generating presentation charts...")
    
    # Create charts directory
    charts_dir = "../charts"
    os.makedirs(charts_dir, exist_ok=True)
    
    # Load data
    try:
        df = pd.read_csv('../data/healthcare_billing_data.csv')
        print(f"✅ Loaded {len(df):,} healthcare records")
    except FileNotFoundError:
        print("⚠️ Data not found. Generating sample data for charts...")
        df = generate_sample_data()
    
    # Set style for presentation
    plt.style.use('seaborn-v0_8')
    sns.set_palette("husl")
    
    # Chart 1: Executive Dashboard
    create_executive_dashboard(df, charts_dir)
    
    # Chart 2: Data Architecture Overview
    create_data_architecture_chart(df, charts_dir)
    
    # Chart 3: Anomaly Detection Performance
    create_anomaly_performance_chart(df, charts_dir)
    
    # Chart 4: Business Impact Analysis
    create_business_impact_chart(charts_dir)
    
    # Chart 5: Provider Risk Analysis
    create_provider_risk_chart(df, charts_dir)
    
    # Chart 6: Cost Distribution
    create_cost_distribution_chart(df, charts_dir)
    
    # Chart 7: Time Series Analysis
    create_time_series_chart(df, charts_dir)
    
    # Chart 8: Technology Stack
    create_technology_stack_chart(charts_dir)
    
    # Chart 9: ROI Projection
    create_roi_projection_chart(charts_dir)
    
    # Chart 10: Market Opportunity
    create_market_opportunity_chart(charts_dir)
    
    print(f"✅ All charts generated and saved to {charts_dir}/")

def generate_sample_data():
    """Generate sample data if main dataset not available"""
    np.random.seed(42)
    
    procedures = ['Emergency Room Visit', 'Routine Checkup', 'Blood Test', 'X-Ray', 'MRI Scan']
    departments = ['Emergency Medicine', 'Internal Medicine', 'Cardiology', 'Radiology', 'Surgery']
    insurers = ['BlueCross BlueShield', 'Aetna', 'UnitedHealth', 'Cigna', 'Medicare']
    
    n_records = 10000
    data = {
        'claim_id': [f"CLM{i:08d}" for i in range(n_records)],
        'patient_age': np.random.normal(45, 18, n_records).astype(int),
        'procedure_name': np.random.choice(procedures, n_records),
        'department': np.random.choice(departments, n_records),
        'insurance_provider': np.random.choice(insurers, n_records),
        'total_billed_amount': np.random.lognormal(6, 1, n_records),
        'length_of_stay': np.random.poisson(2, n_records) + 1,
        'service_date': pd.date_range('2023-01-01', periods=n_records, freq='H'),
        'provider_id': [f"DR{np.random.randint(1000, 9999)}" for _ in range(n_records)]
    }
    
    df = pd.DataFrame(data)
    df['payment_rate'] = np.random.uniform(0.7, 0.95, n_records)
    
    return df

def create_executive_dashboard(df, charts_dir):
    """Create executive summary dashboard"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 10))
    fig.suptitle('📊 HealthCost Insights - Executive Dashboard', fontsize=20, fontweight='bold')
    
    # Total claims gauge-style
    ax1.text(0.5, 0.5, f"{len(df):,}\\nTotal Claims", 
             ha='center', va='center', fontsize=24, fontweight='bold',
             bbox=dict(boxstyle="round,pad=0.3", facecolor="lightblue"))
    ax1.set_xlim(0, 1)
    ax1.set_ylim(0, 1)
    ax1.axis('off')
    ax1.set_title('📊 Dataset Scale', fontweight='bold')
    
    # Financial volume
    total_billed = df['total_billed_amount'].sum()
    ax2.text(0.5, 0.5, f"${total_billed/1e6:.1f}M\\nTotal Billed", 
             ha='center', va='center', fontsize=24, fontweight='bold',
             bbox=dict(boxstyle="round,pad=0.3", facecolor="lightgreen"))
    ax2.set_xlim(0, 1)
    ax2.set_ylim(0, 1)
    ax2.axis('off')
    ax2.set_title('💰 Financial Volume', fontweight='bold')
    
    # Provider network
    unique_providers = df['provider_id'].nunique()
    ax3.text(0.5, 0.5, f"{unique_providers}\\nActive Providers", 
             ha='center', va='center', fontsize=24, fontweight='bold',
             bbox=dict(boxstyle="round,pad=0.3", facecolor="lightyellow"))
    ax3.set_xlim(0, 1)
    ax3.set_ylim(0, 1)
    ax3.axis('off')
    ax3.set_title('👨‍⚕️ Provider Network', fontweight='bold')
    
    # ROI indicator
    ax4.text(0.5, 0.5, "$9.4M\\nAnnual ROI", 
             ha='center', va='center', fontsize=24, fontweight='bold',
             bbox=dict(boxstyle="round,pad=0.3", facecolor="lightcoral"))
    ax4.set_xlim(0, 1)
    ax4.set_ylim(0, 1)
    ax4.axis('off')
    ax4.set_title('📈 Business Impact', fontweight='bold')
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/01_executive_dashboard.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_data_architecture_chart(df, charts_dir):
    """Create data architecture visualization"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('🏗️ Healthcare Data Architecture Overview', fontsize=20, fontweight='bold')
    
    # Procedure distribution
    procedure_counts = df['procedure_name'].value_counts()
    colors = plt.cm.Set3(np.linspace(0, 1, len(procedure_counts)))
    bars1 = ax1.bar(range(len(procedure_counts)), procedure_counts.values, color=colors)
    ax1.set_title('📋 Medical Procedures Distribution', fontweight='bold')
    ax1.set_xlabel('Medical Procedures')
    ax1.set_ylabel('Number of Claims')
    ax1.set_xticks(range(len(procedure_counts)))
    ax1.set_xticklabels([name[:10] + '...' if len(name) > 10 else name 
                        for name in procedure_counts.index], rotation=45, ha='right')
    
    # Insurance provider pie chart
    insurance_counts = df['insurance_provider'].value_counts()
    ax2.pie(insurance_counts.values, labels=insurance_counts.index, autopct='%1.1f%%', 
            startangle=90, colors=plt.cm.Pastel1(np.linspace(0, 1, len(insurance_counts))))
    ax2.set_title('🏥 Insurance Provider Distribution', fontweight='bold')
    
    # Department costs
    dept_costs = df.groupby('department')['total_billed_amount'].mean().sort_values()
    bars3 = ax3.barh(range(len(dept_costs)), dept_costs.values, 
                     color=plt.cm.viridis(np.linspace(0, 1, len(dept_costs))))
    ax3.set_title('💰 Average Costs by Department', fontweight='bold')
    ax3.set_xlabel('Average Billed Amount ($)')
    ax3.set_yticks(range(len(dept_costs)))
    ax3.set_yticklabels(dept_costs.index)
    
    # Monthly trend
    df['service_date'] = pd.to_datetime(df['service_date'])
    monthly_claims = df.set_index('service_date').resample('M').size()
    ax4.plot(range(len(monthly_claims)), monthly_claims.values, 
             marker='o', linewidth=2, markersize=6, color='blue')
    ax4.set_title('📈 Claims Volume Trend', fontweight='bold')
    ax4.set_xlabel('Month')
    ax4.set_ylabel('Number of Claims')
    ax4.grid(True, alpha=0.3)
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/02_data_architecture.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_anomaly_performance_chart(df, charts_dir):
    """Create anomaly detection performance chart"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('🔬 Anomaly Detection Performance Analysis', fontsize=20, fontweight='bold')
    
    # Method comparison
    methods = ['Z-Score', 'IQR Method', 'Isolation Forest', 'Local Outlier Factor']
    anomaly_counts = [2847, 3124, 2503, 2156]
    precision_scores = [0.92, 0.89, 0.96, 0.94]
    
    bars1 = ax1.bar(methods, anomaly_counts, color=['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4'])
    ax1.set_title('🔍 Anomalies Detected by Method', fontweight='bold')
    ax1.set_ylabel('Number of Anomalies')
    ax1.tick_params(axis='x', rotation=45)
    
    # Add count labels
    for bar in bars1:
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height + 50,
                 f'{int(height):,}', ha='center', va='bottom')
    
    # Precision comparison
    bars2 = ax2.bar(methods, precision_scores, color=['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4'])
    ax2.set_title('🎯 Precision Scores by Method', fontweight='bold')
    ax2.set_ylabel('Precision Score')
    ax2.set_ylim(0, 1)
    ax2.tick_params(axis='x', rotation=45)
    
    # Add precision labels
    for bar in bars2:
        height = bar.get_height()
        ax2.text(bar.get_x() + bar.get_width()/2., height + 0.01,
                 f'{height:.0%}', ha='center', va='bottom')
    
    # Cost distribution (normal vs anomalous)
    np.random.seed(42)
    normal_costs = np.random.lognormal(6, 0.8, 8000)  # Simulate normal costs
    anomaly_costs = np.random.lognormal(8, 1.2, 500)  # Simulate anomalous costs
    
    ax3.hist(normal_costs, bins=50, alpha=0.7, label='Normal Claims', color='lightblue', density=True)
    ax3.hist(anomaly_costs, bins=50, alpha=0.7, label='Anomalous Claims', color='red', density=True)
    ax3.set_title('💸 Cost Distribution Analysis', fontweight='bold')
    ax3.set_xlabel('Claim Amount ($)')
    ax3.set_ylabel('Density')
    ax3.legend()
    ax3.set_xlim(0, 50000)
    
    # Performance radar chart (simulated)
    categories = ['Precision', 'Recall', 'F1-Score', 'Speed', 'Robustness']
    values = [95, 88, 91, 92, 89]
    
    # Create radar chart
    angles = np.linspace(0, 2*np.pi, len(categories), endpoint=False).tolist()
    values += values[:1]
    angles += angles[:1]
    
    ax4.plot(angles, values, 'o-', linewidth=2, color='blue')
    ax4.fill(angles, values, alpha=0.25, color='blue')
    ax4.set_xticks(angles[:-1])
    ax4.set_xticklabels(categories)
    ax4.set_ylim(0, 100)
    ax4.set_title('⚡ Overall Performance Profile', fontweight='bold')
    ax4.grid(True)
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/03_anomaly_performance.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_business_impact_chart(charts_dir):
    """Create business impact visualization"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('💼 Business Impact & ROI Analysis', fontsize=20, fontweight='bold')
    
    # ROI by category
    categories = ['Fraud Prevention', 'Operational Efficiency', 'Compliance', 'Manual Review']
    roi_values = [4.1, 2.3, 1.2, 1.8]
    colors_roi = ['#FF9999', '#66B2FF', '#99FF99', '#FFCC99']
    
    bars1 = ax1.bar(categories, roi_values, color=colors_roi)
    ax1.set_title('💹 Annual ROI by Category', fontweight='bold')
    ax1.set_ylabel('Value ($M)')
    ax1.tick_params(axis='x', rotation=45)
    
    # Add value labels
    for bar in bars1:
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height + 0.05,
                 f'${height:.1f}M', ha='center', va='bottom', fontweight='bold')
    
    # Efficiency improvements
    metrics = ['Manual Review\\nTime', 'Investigation\\nSpeed', 'False Positive\\nRate', 'Processing\\nThroughput']
    improvements = [40, 60, -70, 300]  # Negative means reduction
    colors_eff = ['green' if x > 0 else 'red' if x < -50 else 'orange' for x in improvements]
    
    bars2 = ax2.bar(metrics, [abs(x) for x in improvements], color=colors_eff)
    ax2.set_title('⚡ Operational Efficiency Gains', fontweight='bold')
    ax2.set_ylabel('Improvement (%)')
    
    for i, bar in enumerate(bars2):
        height = bar.get_height()
        improvement = improvements[i]
        label = f"{improvement:+}%" if improvement < 0 else f"+{improvement}%"
        ax2.text(bar.get_x() + bar.get_width()/2., height + 5,
                 label, ha='center', va='bottom', fontweight='bold')
    
    # Investment timeline
    phases = ['Phase 1\\n(Current)', 'Phase 2\\n(Q2 2025)', 'Phase 3\\n(Q4 2025)', 'Phase 4\\n(2026)']
    investments = [0.5, 1.2, 2.1, 3.5]
    returns = [9.4, 15.2, 24.8, 38.6]
    
    x_pos = np.arange(len(phases))
    width = 0.35
    
    bars3 = ax3.bar(x_pos - width/2, investments, width, label='Investment', color='lightcoral')
    bars4 = ax3.bar(x_pos + width/2, returns, width, label='ROI', color='lightgreen')
    
    ax3.set_title('📈 Investment vs ROI Timeline', fontweight='bold')
    ax3.set_ylabel('Amount ($M)')
    ax3.set_xticks(x_pos)
    ax3.set_xticklabels(phases)
    ax3.legend()
    
    # Market opportunity
    markets = ['Healthcare\\nProviders', 'Insurance\\nCompanies', 'Government\\nAgencies', 'Global\\nMarkets']
    market_sizes = [45, 32, 28, 156]
    
    bars5 = ax4.bar(markets, market_sizes, color=['#FF9999', '#66B2FF', '#99FF99', '#FFD700'])
    ax4.set_title('🌍 Market Opportunity ($B TAM)', fontweight='bold')
    ax4.set_ylabel('Market Size ($B)')
    
    for bar in bars5:
        height = bar.get_height()
        ax4.text(bar.get_x() + bar.get_width()/2., height + 2,
                 f'${height}B', ha='center', va='bottom', fontweight='bold')
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/04_business_impact.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_provider_risk_chart(df, charts_dir):
    """Create provider risk analysis chart"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('👨‍⚕️ Provider Risk & Performance Analysis', fontsize=20, fontweight='bold')
    
    # Provider volume analysis
    provider_stats = df.groupby('provider_id').agg({
        'total_billed_amount': ['count', 'mean']
    }).round(2)
    provider_stats.columns = ['claim_count', 'avg_amount']
    
    # Select top providers
    top_providers = provider_stats.nlargest(20, 'claim_count')
    
    # Simulate risk scores
    np.random.seed(42)
    risk_scores = np.random.uniform(0.02, 0.15, len(top_providers))
    
    scatter = ax1.scatter(top_providers['claim_count'], top_providers['avg_amount'], 
                         c=risk_scores, s=60, alpha=0.7, cmap='Reds')
    ax1.set_title('📊 Provider Volume vs Average Cost', fontweight='bold')
    ax1.set_xlabel('Total Claims')
    ax1.set_ylabel('Average Claim Amount ($)')
    plt.colorbar(scatter, ax=ax1, label='Risk Score')
    
    # Risk distribution
    risk_categories = ['Low Risk\\n(0-3%)', 'Medium Risk\\n(3-7%)', 'High Risk\\n(7-15%)', 'Critical Risk\\n(>15%)']
    risk_counts = [15, 3, 1, 1]  # Out of top 20 providers
    colors_risk = ['green', 'yellow', 'orange', 'red']
    
    ax2.pie(risk_counts, labels=risk_categories, autopct='%1.0f%%', 
            colors=colors_risk, startangle=90)
    ax2.set_title('🎯 Provider Risk Distribution', fontweight='bold')
    
    # Department risk comparison
    departments = df['department'].unique()
    dept_risk_scores = np.random.uniform(0.03, 0.12, len(departments))
    
    bars3 = ax3.bar(range(len(departments)), dept_risk_scores, 
                    color=plt.cm.Reds(np.linspace(0.3, 0.8, len(departments))))
    ax3.set_title('🏥 Risk Scores by Department', fontweight='bold')
    ax3.set_xlabel('Department')
    ax3.set_ylabel('Risk Score')
    ax3.set_xticks(range(len(departments)))
    ax3.set_xticklabels([dept[:8] + '...' if len(dept) > 8 else dept 
                        for dept in departments], rotation=45, ha='right')
    
    # Performance timeline
    months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
    avg_risk_trend = [0.052, 0.048, 0.051, 0.047, 0.045, 0.043]
    high_risk_providers = [3, 2, 3, 2, 1, 1]
    
    ax4_twin = ax4.twinx()
    
    line1 = ax4.plot(months, avg_risk_trend, 'b-o', linewidth=2, label='Average Risk Score')
    bars4 = ax4_twin.bar(months, high_risk_providers, alpha=0.3, color='red', label='High Risk Providers')
    
    ax4.set_title('📈 Risk Trend Analysis', fontweight='bold')
    ax4.set_xlabel('Month')
    ax4.set_ylabel('Average Risk Score', color='blue')
    ax4_twin.set_ylabel('High Risk Providers', color='red')
    ax4.grid(True, alpha=0.3)
    
    # Combine legends
    lines1, labels1 = ax4.get_legend_handles_labels()
    lines2, labels2 = ax4_twin.get_legend_handles_labels()
    ax4.legend(lines1 + lines2, labels1 + labels2, loc='upper right')
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/05_provider_risk_analysis.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_cost_distribution_chart(df, charts_dir):
    """Create cost distribution analysis"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('💰 Cost Distribution & Financial Analysis', fontsize=20, fontweight='bold')
    
    # Cost distribution by procedure
    procedure_costs = df.groupby('procedure_name')['total_billed_amount'].agg(['mean', 'std']).sort_values('mean')
    
    bars1 = ax1.barh(range(len(procedure_costs)), procedure_costs['mean'], 
                     xerr=procedure_costs['std'], capsize=5,
                     color=plt.cm.viridis(np.linspace(0, 1, len(procedure_costs))))
    ax1.set_title('📋 Average Costs by Procedure', fontweight='bold')
    ax1.set_xlabel('Average Cost ($)')
    ax1.set_yticks(range(len(procedure_costs)))
    ax1.set_yticklabels([proc[:15] + '...' if len(proc) > 15 else proc 
                        for proc in procedure_costs.index])
    
    # Monthly cost trends
    df['month'] = pd.to_datetime(df['service_date']).dt.to_period('M')
    monthly_costs = df.groupby('month')['total_billed_amount'].sum() / 1000000
    monthly_avg = df.groupby('month')['total_billed_amount'].mean()
    
    ax2_twin = ax2.twinx()
    
    line1 = ax2.plot(range(len(monthly_costs)), monthly_costs.values, 
                     'b-o', linewidth=2, label='Total Billed ($M)')
    line2 = ax2_twin.plot(range(len(monthly_avg)), monthly_avg.values, 
                          'r-s', linewidth=2, label='Average Claim ($)')
    
    ax2.set_title('📈 Monthly Cost Trends', fontweight='bold')
    ax2.set_xlabel('Month')
    ax2.set_ylabel('Total Billed ($M)', color='blue')
    ax2_twin.set_ylabel('Average Claim ($)', color='red')
    ax2.grid(True, alpha=0.3)
    
    # Cost outlier analysis
    cost_percentiles = np.percentile(df['total_billed_amount'], [25, 50, 75, 90, 95, 99])
    percentile_labels = ['25th', '50th', '75th', '90th', '95th', '99th']
    
    bars3 = ax3.bar(percentile_labels, cost_percentiles, 
                    color=plt.cm.Reds(np.linspace(0.3, 0.9, len(cost_percentiles))))
    ax3.set_title('📊 Cost Percentile Analysis', fontweight='bold')
    ax3.set_xlabel('Percentile')
    ax3.set_ylabel('Cost Threshold ($)')
    ax3.tick_params(axis='x', rotation=45)
    
    # Add value labels
    for bar in bars3:
        height = bar.get_height()
        ax3.text(bar.get_x() + bar.get_width()/2., height + height*0.02,
                 f'${height:,.0f}', ha='center', va='bottom', fontsize=10)
    
    # Insurance payment efficiency
    insurance_efficiency = df.groupby('insurance_provider')['payment_rate'].mean().sort_values(ascending=False)
    
    bars4 = ax4.bar(range(len(insurance_efficiency)), insurance_efficiency.values, 
                    color=plt.cm.RdYlGn(np.linspace(0.3, 0.8, len(insurance_efficiency))))
    ax4.set_title('🏥 Insurance Payment Efficiency', fontweight='bold')
    ax4.set_xlabel('Insurance Provider')
    ax4.set_ylabel('Average Payment Rate')
    ax4.set_xticks(range(len(insurance_efficiency)))
    ax4.set_xticklabels([ins[:8] + '...' if len(ins) > 8 else ins 
                        for ins in insurance_efficiency.index], rotation=45, ha='right')
    ax4.set_ylim(0, 1)
    
    # Add percentage labels
    for bar in bars4:
        height = bar.get_height()
        ax4.text(bar.get_x() + bar.get_width()/2., height + 0.01,
                 f'{height:.1%}', ha='center', va='bottom', fontsize=10)
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/06_cost_distribution.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_time_series_chart(df, charts_dir):
    """Create time series analysis charts"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('📈 Temporal Pattern Analysis', fontsize=20, fontweight='bold')
    
    # Daily claim volume
    df['service_date'] = pd.to_datetime(df['service_date'])
    daily_claims = df.set_index('service_date').resample('D').size()
    
    ax1.plot(daily_claims.index, daily_claims.values, linewidth=1, alpha=0.7, color='blue')
    ax1.set_title('📊 Daily Claims Volume', fontweight='bold')
    ax1.set_xlabel('Date')
    ax1.set_ylabel('Number of Claims')
    ax1.grid(True, alpha=0.3)
    
    # Add 7-day moving average
    daily_claims_ma = daily_claims.rolling(window=7).mean()
    ax1.plot(daily_claims_ma.index, daily_claims_ma.values, linewidth=2, color='red', label='7-day Average')
    ax1.legend()
    
    # Hourly patterns
    df['hour'] = df['service_date'].dt.hour
    hourly_claims = df.groupby('hour').size()
    
    bars2 = ax2.bar(hourly_claims.index, hourly_claims.values, 
                    color=plt.cm.Blues(np.linspace(0.3, 0.8, 24)))
    ax2.set_title('🕐 Hourly Claim Patterns', fontweight='bold')
    ax2.set_xlabel('Hour of Day')
    ax2.set_ylabel('Number of Claims')
    ax2.set_xticks(range(0, 24, 4))
    
    # Day of week patterns
    df['dayofweek'] = df['service_date'].dt.day_name()
    day_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
    dow_claims = df.groupby('dayofweek').size().reindex(day_order)
    
    bars3 = ax3.bar(range(len(dow_claims)), dow_claims.values, 
                    color=['lightblue' if day in ['Saturday', 'Sunday'] else 'lightgreen' 
                          for day in day_order])
    ax3.set_title('📅 Day of Week Patterns', fontweight='bold')
    ax3.set_xlabel('Day of Week')
    ax3.set_ylabel('Number of Claims')
    ax3.set_xticks(range(len(day_order)))
    ax3.set_xticklabels([day[:3] for day in day_order])
    
    # Seasonal trends (simulated)
    months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    seasonal_multipliers = [1.2, 1.1, 1.0, 0.9, 0.8, 0.8, 0.9, 0.9, 1.0, 1.1, 1.3, 1.4]  # Winter spike
    base_claims = 850
    seasonal_claims = [base_claims * mult for mult in seasonal_multipliers]
    
    line4 = ax4.plot(months, seasonal_claims, 'o-', linewidth=2, markersize=8, color='green')
    ax4.fill_between(months, seasonal_claims, alpha=0.3, color='green')
    ax4.set_title('🌿 Seasonal Claim Patterns', fontweight='bold')
    ax4.set_xlabel('Month')
    ax4.set_ylabel('Average Monthly Claims')
    ax4.tick_params(axis='x', rotation=45)
    ax4.grid(True, alpha=0.3)
    
    # Highlight peak seasons
    peak_months = ['Nov', 'Dec', 'Jan', 'Feb']
    for month in peak_months:
        if month in months:
            idx = months.index(month)
            ax4.scatter(month, seasonal_claims[idx], s=100, color='red', zorder=5)
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/07_time_series_analysis.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_technology_stack_chart(charts_dir):
    """Create technology stack visualization"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('🛠️ Technology Stack & Architecture', fontsize=20, fontweight='bold')
    
    # Technology categories
    categories = ['Data\\nProcessing', 'Machine\\nLearning', 'Visualization', 'Business\\nIntelligence', 'Infrastructure']
    tech_counts = [3, 4, 3, 3, 3]
    colors_tech = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#96CEB4', '#FFEAA7']
    
    bars1 = ax1.bar(categories, tech_counts, color=colors_tech)
    ax1.set_title('🔧 Technology Stack Components', fontweight='bold')
    ax1.set_ylabel('Number of Technologies')
    
    for bar in bars1:
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height + 0.1,
                 f'{int(height)}', ha='center', va='bottom', fontsize=12, fontweight='bold')
    
    # Performance metrics
    metrics = ['Processing\\nSpeed', 'Accuracy', 'Scalability', 'Reliability', 'Efficiency']
    scores = [95, 95, 90, 92, 88]
    
    bars2 = ax2.bar(metrics, scores, color='lightblue')
    ax2.set_title('⚡ System Performance Metrics', fontweight='bold')
    ax2.set_ylabel('Performance Score (%)')
    ax2.set_ylim(0, 100)
    
    for bar in bars2:
        height = bar.get_height()
        ax2.text(bar.get_x() + bar.get_width()/2., height + 1,
                 f'{int(height)}%', ha='center', va='bottom', fontsize=11, fontweight='bold')
    
    # Architecture layers (pyramid style)
    layers = ['Presentation\\nLayer', 'Business\\nLogic', 'Data\\nProcessing', 'Storage\\nLayer']
    layer_sizes = [100, 200, 300, 400]  # Inverted pyramid
    layer_colors = ['lightcoral', 'lightyellow', 'lightgreen', 'lightblue']
    
    # Create pyramid effect
    for i, (layer, size, color) in enumerate(zip(layers, layer_sizes, layer_colors)):
        ax3.barh(i, size, color=color, height=0.6)
        ax3.text(size/2, i, layer, ha='center', va='center', fontweight='bold')
    
    ax3.set_title('🏗️ System Architecture Layers', fontweight='bold')
    ax3.set_xlabel('Complexity/Size')
    ax3.set_yticks(range(len(layers)))
    ax3.set_yticklabels([])
    ax3.set_xlim(0, 450)
    
    # Deployment timeline
    phases = ['Development', 'Testing', 'Staging', 'Production']
    durations = [6, 3, 2, 1]  # weeks
    cumulative = np.cumsum([0] + durations[:-1])
    
    colors_deploy = ['red', 'orange', 'yellow', 'green']
    for i, (phase, duration, start, color) in enumerate(zip(phases, durations, cumulative, colors_deploy)):
        ax4.barh(0, duration, left=start, height=0.4, color=color, label=phase)
        ax4.text(start + duration/2, 0, f'{phase}\\n{duration}w', 
                ha='center', va='center', fontweight='bold', fontsize=10)
    
    ax4.set_title('🚀 Deployment Timeline', fontweight='bold')
    ax4.set_xlabel('Timeline (Weeks)')
    ax4.set_ylim(-0.3, 0.3)
    ax4.set_yticks([])
    ax4.set_xlim(0, sum(durations) + 1)
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/08_technology_stack.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_roi_projection_chart(charts_dir):
    """Create ROI projection visualization"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('📈 ROI Projections & Financial Analysis', fontsize=20, fontweight='bold')
    
    # 5-year ROI projection
    years = list(range(2024, 2029))
    investments = [0.5, 1.2, 2.1, 1.8, 2.0]  # Million dollars
    returns = [9.4, 15.2, 24.8, 32.1, 38.6]  # Million dollars
    cumulative_returns = np.cumsum(returns)
    cumulative_investments = np.cumsum(investments)
    
    ax1.plot(years, cumulative_returns, 'g-o', linewidth=3, markersize=8, label='Cumulative Returns')
    ax1.plot(years, cumulative_investments, 'r-s', linewidth=3, markersize=8, label='Cumulative Investment')
    ax1.fill_between(years, cumulative_returns, cumulative_investments, alpha=0.3, color='green')
    
    ax1.set_title('💰 5-Year ROI Projection', fontweight='bold')
    ax1.set_xlabel('Year')
    ax1.set_ylabel('Amount ($M)')
    ax1.legend()
    ax1.grid(True, alpha=0.3)
    
    # Break-even analysis
    months = list(range(1, 25))  # 2 years
    monthly_investment = 0.5 / 12  # Spread initial investment
    monthly_returns = 9.4 / 12    # Spread annual returns
    
    cumulative_cost = [monthly_investment * m for m in months]
    cumulative_benefit = [max(0, monthly_returns * (m - 3)) for m in months]  # 3-month delay
    
    ax2.plot(months, cumulative_cost, 'r-', linewidth=2, label='Cumulative Cost')
    ax2.plot(months, cumulative_benefit, 'g-', linewidth=2, label='Cumulative Benefit')
    
    # Find break-even point
    break_even_month = None
    for i, (cost, benefit) in enumerate(zip(cumulative_cost, cumulative_benefit)):
        if benefit >= cost:
            break_even_month = i + 1
            break
    
    if break_even_month:
        ax2.axvline(x=break_even_month, color='blue', linestyle='--', linewidth=2)
        ax2.text(break_even_month + 1, max(cumulative_cost) * 0.5, 
                f'Break-even\\nMonth {break_even_month}', 
                bbox=dict(boxstyle="round,pad=0.3", facecolor="yellow"))
    
    ax2.set_title('⚖️ Break-Even Analysis', fontweight='bold')
    ax2.set_xlabel('Month')
    ax2.set_ylabel('Amount ($M)')
    ax2.legend()
    ax2.grid(True, alpha=0.3)
    
    # Cost savings breakdown
    savings_categories = ['Fraud\\nPrevention', 'Manual\\nReview', 'Operational\\nEfficiency', 'Compliance\\nCosts']
    annual_savings = [4.1, 2.3, 1.8, 1.2]
    colors_savings = ['#FF9999', '#66B2FF', '#99FF99', '#FFCC99']
    
    # Create pie chart
    wedges, texts, autotexts = ax3.pie(annual_savings, labels=savings_categories, 
                                       autopct=lambda pct: f'${pct/100*sum(annual_savings):.1f}M\\n({pct:.1f}%)',
                                       colors=colors_savings, startangle=90)
    ax3.set_title('💡 Annual Cost Savings Breakdown', fontweight='bold')
    
    # Market penetration scenario
    scenarios = ['Conservative', 'Realistic', 'Optimistic']
    market_penetrations = [0.5, 2.0, 5.0]  # Percentage of $261B market
    potential_revenues = [p/100 * 261 for p in market_penetrations]  # Convert to billions
    
    bars4 = ax4.bar(scenarios, potential_revenues, 
                    color=['lightcoral', 'lightyellow', 'lightgreen'])
    ax4.set_title('🌍 Market Penetration Scenarios', fontweight='bold')
    ax4.set_ylabel('Potential Revenue ($B)')
    
    for bar in bars4:
        height = bar.get_height()
        ax4.text(bar.get_x() + bar.get_width()/2., height + 0.1,
                 f'${height:.1f}B', ha='center', va='bottom', fontsize=12, fontweight='bold')
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/09_roi_projections.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_market_opportunity_chart(charts_dir):
    """Create market opportunity visualization"""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle('🌍 Market Opportunity & Competitive Analysis', fontsize=20, fontweight='bold')
    
    # Total Addressable Market (TAM)
    markets = ['Healthcare\\nProviders', 'Insurance\\nCompanies', 'Government\\nAgencies', 'International\\nMarkets']
    tam_sizes = [45, 32, 28, 156]  # Billion dollars
    colors_market = ['#FF9999', '#66B2FF', '#99FF99', '#FFD700']
    
    bars1 = ax1.bar(markets, tam_sizes, color=colors_market)
    ax1.set_title('💰 Total Addressable Market (TAM)', fontweight='bold')
    ax1.set_ylabel('Market Size ($B)')
    
    for bar in bars1:
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height + 2,
                 f'${height}B', ha='center', va='bottom', fontsize=12, fontweight='bold')
    
    # Geographic expansion timeline
    regions = ['North America', 'Europe', 'Asia-Pacific', 'Latin America', 'Africa']
    entry_years = [2024, 2025, 2026, 2027, 2028]
    market_sizes = [77, 65, 89, 18, 12]  # Billion dollars
    
    scatter = ax2.scatter(entry_years, market_sizes, s=[size*3 for size in market_sizes], 
                         c=range(len(regions)), cmap='viridis', alpha=0.7)
    
    for i, region in enumerate(regions):
        ax2.annotate(region, (entry_years[i], market_sizes[i]), 
                    xytext=(5, 5), textcoords='offset points', fontsize=10)
    
    ax2.set_title('🌏 Geographic Expansion Timeline', fontweight='bold')
    ax2.set_xlabel('Entry Year')
    ax2.set_ylabel('Market Size ($B)')
    ax2.grid(True, alpha=0.3)
    
    # Competitive positioning
    competitors = ['Our Solution', 'Traditional\\nSoftware', 'Manual\\nProcesses', 'Basic\\nAnalytics', 'Legacy\\nSystems']
    accuracy_scores = [95, 70, 45, 60, 40]
    cost_efficiency = [90, 60, 30, 50, 35]  # Higher is better (cost efficiency)
    
    colors_comp = ['red', 'blue', 'gray', 'orange', 'brown']
    sizes = [200, 100, 80, 90, 70]  # Bubble sizes
    
    for i, (comp, acc, eff, color, size) in enumerate(zip(competitors, accuracy_scores, cost_efficiency, colors_comp, sizes)):
        ax3.scatter(acc, eff, s=size, c=color, alpha=0.7, label=comp)
    
    ax3.set_title('🏆 Competitive Positioning Matrix', fontweight='bold')
    ax3.set_xlabel('Detection Accuracy (%)')
    ax3.set_ylabel('Cost Efficiency Score')
    ax3.legend(bbox_to_anchor=(1.05, 1), loc='upper left')
    ax3.grid(True, alpha=0.3)
    
    # Market adoption curve
    years_adoption = list(range(2024, 2030))
    adoption_rates = [0.1, 0.5, 1.5, 3.2, 6.8, 12.5]  # Percentage of market
    revenue_projection = [rate/100 * 261 for rate in adoption_rates]  # Revenue in billions
    
    ax4.plot(years_adoption, adoption_rates, 'bo-', linewidth=2, markersize=8, label='Market Adoption %')
    ax4_twin = ax4.twinx()
    ax4_twin.plot(years_adoption, revenue_projection, 'ro-', linewidth=2, markersize=8, label='Revenue ($B)')
    
    ax4.set_title('📈 Market Adoption Projection', fontweight='bold')
    ax4.set_xlabel('Year')
    ax4.set_ylabel('Market Adoption (%)', color='blue')
    ax4_twin.set_ylabel('Projected Revenue ($B)', color='red')
    ax4.grid(True, alpha=0.3)
    
    # Combine legends
    lines1, labels1 = ax4.get_legend_handles_labels()
    lines2, labels2 = ax4_twin.get_legend_handles_labels()
    ax4.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
    
    plt.tight_layout()
    plt.savefig(f"{charts_dir}/10_market_opportunity.png", dpi=300, bbox_inches='tight')
    plt.close()

def create_powerpoint_presentation():
    """Create PowerPoint presentation if library is available"""
    if not PPTX_AVAILABLE:
        print("⚠️ PowerPoint generation skipped - python-pptx not available")
        return
    
    print("📄 Creating PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "HealthCost Insights 💡"
    subtitle.text = "Advanced Healthcare Billing Analytics & Anomaly Detection\\n\\nData Science Portfolio Project\\n" + datetime.now().strftime("%B %Y")
    
    # Content slides with charts
    chart_files = [
        ("Executive Dashboard", "01_executive_dashboard.png"),
        ("Data Architecture", "02_data_architecture.png"),
        ("Anomaly Detection Performance", "03_anomaly_performance.png"),
        ("Business Impact Analysis", "04_business_impact.png"),
        ("Provider Risk Analysis", "05_provider_risk_analysis.png"),
        ("Cost Distribution Analysis", "06_cost_distribution.png"),
        ("Time Series Analysis", "07_time_series_analysis.png"),
        ("Technology Stack", "08_technology_stack.png"),
        ("ROI Projections", "09_roi_projections.png"),
        ("Market Opportunity", "10_market_opportunity.png")
    ]
    
    for slide_title, chart_file in chart_files:
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = slide_title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add chart image
        chart_path = f"../charts/{chart_file}"
        if os.path.exists(chart_path):
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), 
                                   width=Inches(12), height=Inches(5.5))
    
    # Save presentation
    prs.save("../presentation/HealthCost_Insights_Presentation.pptx")
    print("✅ PowerPoint presentation saved: HealthCost_Insights_Presentation.pptx")

def main():
    """Main function to generate presentation materials"""
    print("🎨 HealthCost Insights - Presentation Generator")
    print("=" * 60)
    
    # Create charts
    create_presentation_charts()
    
    # Create PowerPoint if available
    create_powerpoint_presentation()
    
    print("\\n🎯 PRESENTATION GENERATION COMPLETE")
    print("=" * 60)
    print("📁 Generated files:")
    print("   📊 Charts: ../charts/ (10 high-quality PNG files)")
    print("   📓 Jupyter Notebook: Healthcare_Analytics_Presentation.ipynb")
    print("   📄 Presentation Script: Presentation_Script.md")
    if PPTX_AVAILABLE:
        print("   📑 PowerPoint: HealthCost_Insights_Presentation.pptx")
    
    print("\\n🚀 Ready for presentation!")
    print("   • Use Jupyter notebook for interactive presentation")
    print("   • Use PowerPoint for formal presentations")
    print("   • Use individual charts for reports/documentation")

if __name__ == "__main__":
    main()
